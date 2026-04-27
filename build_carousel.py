# “””
RELEASED Healing Method — Carousel PPTX Builder

Flask service. Make POSTs Claude’s JSON output to /build-carousel.
Returns a .pptx file as a binary response.

Deploy to Railway, Render, or any Python host.
Set PHOTOS_DIR environment variable to the folder containing PHOTO-A.jpg through PHOTO-G.jpg.

Dependencies: flask, python-pptx, Pillow, numpy, lxml
“””

import io
import json
import os
import sys
import tempfile
from copy import deepcopy

import numpy as np
from flask import Flask, request, send_file, jsonify
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu
from lxml import etree

app = Flask(**name**)

# ─────────────────────────────────────────────

# CANVAS

# ─────────────────────────────────────────────

SLIDE_W = Inches(11.25)
SLIDE_H = Inches(14.0625)
LM = Inches(0.9375) # left margin
TW = Inches(9.375) # text width
TOP_S1 = Inches(2.25) # headline top — slide 1
TOP_INT = Inches(0.875) # section label top — interior slides

# ─────────────────────────────────────────────

# COLORS

# ─────────────────────────────────────────────

def s(val):
“”“Safely convert any value to string. Returns empty string for None/dict/list.”””
if val is None:
return ‘’
if isinstance(val, (dict, list)):
return ‘’
return str(val)

def rgb(h):
“”“Convert hex string to RGBColor.”””
h = h.lstrip(’#’)
return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def rgb_tuple(h):
h = h.lstrip(’#’)
return (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# Brand palette — exact values only

C = {
‘cream’: ‘#FAF6EE’,
‘offwhite’: ‘#F5F0E8’,
‘dark_brown’: ‘#4a403a’,
‘fg_green’: ‘#11836b’,
‘fg_deep’: ‘#2D3D31’,
‘teal’: ‘#6fafb8’,
‘terra’: ‘#c47a6a’,
‘terra_deep’: ‘#a0503c’,
‘gold’: ‘#e6c98a’,
‘sage’: ‘#83b7a8’,
‘sage_deep’: ‘#4a7a6e’,
‘salmon’: ‘#e6a78f’,
‘deep_terra’: ‘#7b2d1a’,
‘mid_terra’: ‘#a34e3a’,
‘teal_deep’: ‘#3a6e78’,
‘espresso’: ‘#3D2E24’,
}

# Gradient recipes: name -> list of hex stops (top to bottom)

GRADIENTS = {
‘deep_terra’: [’#7b2d1a’, ‘#a34e3a’, ‘#c47a6a’],
‘terra_salmon’: [’#c47a6a’, ‘#e6a78f’],
‘green_teal’: [’#11836b’, ‘#6fafb8’],
‘fg_deep’: [’#2c3d31’, ‘#1a493b’],
‘sage_sagedeep’: [’#4a7a6e’, ‘#83b7a8’],
‘forest_teal’: [’#1c4a38’, ‘#2d7a6a’],
}

# ─────────────────────────────────────────────

# GRADIENT GENERATION (PIL / numpy)

# ─────────────────────────────────────────────

def make_gradient_image(stops, width=1080, height=1350):
“””
Build a vertical gradient from a list of hex color stops.
Uses numpy interpolation — never PPTX native gradients.
“””
stops_rgb = [rgb_tuple(s) for s in stops]
n = len(stops_rgb)
arr = np.zeros((height, width, 3), dtype=np.uint8)

```
segment_h = height // (n - 1)
for seg in range(n - 1):
c1 = stops_rgb[seg]
c2 = stops_rgb[seg + 1]
y_start = seg * segment_h
y_end = (seg + 1) * segment_h if seg < n - 2 else height
seg_len = y_end - y_start
for i in range(seg_len):
t = i / max(seg_len - 1, 1)
row = y_start + i
arr[row, :] = [
int(c1[j] + t * (c2[j] - c1[j])) for j in range(3)
]
return Image.fromarray(arr, 'RGB')
```

def add_gradient_bg(slide, stops):
“”“Add a PIL gradient as a full-slide background image at z-index 2.”””
img = make_gradient_image(stops)
buf = io.BytesIO()
img.save(buf, format=‘PNG’)
buf.seek(0)
pic = slide.shapes.add_picture(buf, 0, 0, SLIDE_W, SLIDE_H)
# Push to back (behind all text boxes)
sp_tree = slide.shapes._spTree
sp_tree.remove(pic._element)
sp_tree.insert(2, pic._element)

def add_solid_bg(slide, hex_color):
“”“Add a solid color rectangle as background.”””
shape = slide.shapes.add_shape(
1, # MSO_SHAPE_TYPE.RECTANGLE
0, 0, SLIDE_W, SLIDE_H
)
shape.fill.solid()
shape.fill.fore_color.rgb = rgb(hex_color)
shape.line.fill.background()
# Push to back
sp_tree = slide.shapes._spTree
sp_tree.remove(shape._element)
sp_tree.insert(2, shape._element)

# ─────────────────────────────────────────────

# PHOTO HANDLING

# ─────────────────────────────────────────────

PHOTOS_DIR = os.environ.get(‘PHOTOS_DIR’, ‘./photos’)

PHOTO_FILES = {
‘PHOTO-A’: ‘PHOTO-A.jpg’,
‘PHOTO-B’: ‘PHOTO-B.jpg’,
‘PHOTO-C’: ‘PHOTO-C.jpg’,
‘PHOTO-D’: ‘PHOTO-D.jpg’,
‘PHOTO-E’: ‘PHOTO-E.jpg’,
‘PHOTO-F’: ‘PHOTO-F.jpg’,
‘PHOTO-G’: ‘PHOTO-G.jpg’,
}

def add_photo_bg(slide, photo_key, tint_color_hex, tint_opacity=0.52):
“””
Load photo, crop to 1080x1350, apply tint, embed at z-index 2.
Text zone: PHOTO-F = top (clear wall zone). PHOTO-A = lower 40%.
Tint opacity range: 0.45-0.62.
“””
path = os.path.join(PHOTOS_DIR, PHOTO_FILES.get(photo_key, ‘PHOTO-A.jpg’))
if not os.path.exists(path):
# Fallback: solid background in tint color if photo missing
add_solid_bg(slide, tint_color_hex)
return

```
img = Image.open(path).convert('RGB')

# Center-crop to 1080x1350
target_w, target_h = 1080, 1350
w_scale = target_w / img.width
h_scale = target_h / img.height
scale = max(w_scale, h_scale)
new_w = int(img.width * scale)
new_h = int(img.height * scale)
img = img.resize((new_w, new_h), Image.LANCZOS)

# For portrait photos (PHOTO-F): crop from top to preserve face at bottom
if photo_key == 'PHOTO-F':
ty = 0
else:
ty = (new_h - target_h) // 2
tx = (new_w - target_w) // 2
img = img.crop((tx, ty, tx + target_w, ty + target_h))

# Apply tint overlay
tint = Image.new('RGB', (target_w, target_h), rgb_tuple(tint_color_hex))
img = Image.blend(img, tint, alpha=tint_opacity)

buf = io.BytesIO()
img.save(buf, format='PNG')
buf.seek(0)

pic = slide.shapes.add_picture(buf, 0, 0, SLIDE_W, SLIDE_H)
sp_tree = slide.shapes._spTree
sp_tree.remove(pic._element)
sp_tree.insert(2, pic._element)
```

# ─────────────────────────────────────────────

# TEXT BOX HELPERS

# ─────────────────────────────────────────────

def add_textbox(slide, text, x, y, w, h, font_name, pt_size,
color_hex, italic=False, align=PP_ALIGN.LEFT, space_after=0):
“”“Add a single-run text box. Returns the shape.”””
txBox = slide.shapes.add_textbox(x, y, w, h)
tf = txBox.text_frame
tf.word_wrap = True
tf.auto_size = None
p = tf.paragraphs[0]
p.alignment = align
run = p.add_run()
run.text = text
run.font.name = font_name
run.font.size = Pt(pt_size)
run.font.color.rgb = rgb(color_hex)
run.font.italic = italic
run.font.bold = False
if space_after:
p.space_after = Pt(space_after)
return txBox

def add_headline_mixed(slide, regular_text, italic_text, x, y, w, h,
font_name, pt_size, reg_color, italic_color,
align=PP_ALIGN.LEFT):
“””
Add a headline with a regular portion and an italic accent portion.
Both on the same line (same paragraph, two runs).
“””
txBox = slide.shapes.add_textbox(x, y, w, h)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = align

```
if regular_text:
r1 = p.add_run()
r1.text = regular_text
r1.font.name = font_name
r1.font.size = Pt(pt_size)
r1.font.color.rgb = rgb(reg_color)
r1.font.italic = False
r1.font.bold = False

if italic_text:
r2 = p.add_run()
r2.text = italic_text
r2.font.name = font_name
r2.font.size = Pt(pt_size)
r2.font.color.rgb = rgb(italic_color)
r2.font.italic = True
r2.font.bold = False

return txBox
```

def set_letter_spacing(run, spc_value=150):
“””
Set letter spacing (character spacing) on a run via XML.
spc=150 in OOXML = 1.5pt letter spacing.
Required for section labels — cannot be set via python-pptx API.
“””
rPr = run._r.get_or_add_rPr()
rPr.set(‘spc’, str(spc_value))

# ─────────────────────────────────────────────

# BRAND ELEMENTS

# ─────────────────────────────────────────────

def add_accent_mark(slide, color_hex):
“””
Short horizontal rule at bottom-left.
Appears on Slide 1 and final slide only. Never interior slides.
“””
shape = slide.shapes.add_shape(
1,
LM, Inches(12.08),
Inches(1.458), Inches(0.104)
)
shape.fill.solid()
shape.fill.fore_color.rgb = rgb(color_hex)
shape.line.fill.background()

def add_section_label(slide, label_text, accent_color_hex):
“””
Section label: short accent bar + tracked-caps Calibri text.
Interior slides 2-5 only. Never Slide 1.
label_text should already be ALL CAPS.
“””
# Accent bar
bar = slide.shapes.add_shape(
1,
LM, Inches(0.90),
Inches(0.43), Inches(0.04)
)
bar.fill.solid()
bar.fill.fore_color.rgb = rgb(accent_color_hex)
bar.line.fill.background()

```
# Label text with letter spacing
txBox = slide.shapes.add_textbox(
Inches(1.47), Inches(0.875),
Inches(8.5), Inches(0.35)
)
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = label_text.upper()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.rgb = rgb(accent_color_hex)
run.font.bold = True
run.font.italic = False
set_letter_spacing(run, 150)
```

# ─────────────────────────────────────────────

# SLIDE BUILDERS

# ─────────────────────────────────────────────

def build_slide_1(prs, copy, carousel_type, photo_assignment, italic_accent_color, gradient_spec):
“””
S1 Hook — type-driven background.
Type 1: full-bleed photo + tint.
Type 2: deep gradient cover.
Type 3: cream background.
No section label. No logo. Accent mark on.
“””
slide_layout = prs.slide_layouts[6] # blank
slide = prs.slides.add_slide(slide_layout)

```
headline = s(copy.get('headline', ''))
headline_it = s(copy.get('headline_italic', ''))
sub = s(copy.get('sub_statement', ''))

if carousel_type == 1:
# Photo + tint — use topic primary dark color as tint
tint_color = C['fg_deep'] if 'green' in gradient_spec.lower() else \
C['terra'] if 'terra' in gradient_spec.lower() else \
C['sage_deep'] if 'sage' in gradient_spec.lower() else C['fg_deep']
add_photo_bg(slide, photo_assignment, tint_color, tint_opacity=0.52)
text_y = Inches(7.5) # lower 40% of slide
hl_color = C['offwhite']
hl_it_color = italic_accent_color
sub_color = C['offwhite']
accent_color = C['offwhite']

elif carousel_type == 2:
# Deep gradient cover
stops = GRADIENTS.get(gradient_spec, GRADIENTS['deep_terra'])
add_gradient_bg(slide, stops)
text_y = TOP_S1
hl_color = C['offwhite']
hl_it_color = italic_accent_color
sub_color = C['offwhite']
accent_color = C['offwhite']

else: # Type 3 — cream
add_solid_bg(slide, C['cream'])
text_y = TOP_S1
hl_color = C['dark_brown']
hl_it_color = italic_accent_color
sub_color = C['dark_brown']
accent_color = C['terra']

# Headline (regular + italic on same slide, stacked as separate paragraphs for line breaks)
if headline and headline_it:
add_headline_mixed(
slide, headline + ' ', headline_it,
LM, text_y, TW, Inches(4.5),
'Georgia', 84, hl_color, hl_it_color
)
elif headline:
add_textbox(slide, headline, LM, text_y, TW, Inches(3.0),
'Georgia', 84, hl_color)

# Sub-statement
if sub:
sub_y = text_y + Inches(5.2) if carousel_type == 1 else Inches(11.0)
add_textbox(slide, sub,
LM, sub_y, TW, Inches(1.0),
'Georgia', 32, sub_color, italic=True)

# Accent mark
add_accent_mark(slide, accent_color)
return slide
```

def build_slide_2(prs, copy, italic_accent_color, is_dark_bg=False, gradient_stops=None):
“””
S2 Validate — cream or warm gradient. Section label. 54pt headline. Body. Italic close.
“””
slide = prs.slides.add_slide(prs.slide_layouts[6])

```
label = s(copy.get('label', ''))
headline = s(copy.get('headline', ''))
hl_it = s(copy.get('headline_italic', ''))
body = s(copy.get('body', ''))
closing = s(copy.get('closing', ''))

if gradient_stops:
add_gradient_bg(slide, gradient_stops)
bg = 'gradient'
elif is_dark_bg:
add_solid_bg(slide, C['fg_deep'])
bg = 'dark'
else:
add_solid_bg(slide, C['cream'])
bg = 'cream'

hl_color = C['offwhite'] if bg in ('dark', 'gradient') else C['dark_brown']
body_color = C['offwhite'] if bg in ('dark', 'gradient') else C['dark_brown']
close_color = C['offwhite'] if bg in ('dark', 'gradient') else italic_accent_color

# Section label
if label:
label_color = C['offwhite'] if bg in ('dark', 'gradient') else italic_accent_color
add_section_label(slide, label, label_color)

# Headline
if headline and hl_it:
add_headline_mixed(slide, headline + ' ', hl_it,
LM, Inches(1.6), TW, Inches(3.5),
'Georgia', 54, hl_color, italic_accent_color if bg == 'cream' else C['offwhite'])
elif headline:
add_textbox(slide, headline, LM, Inches(1.6), TW, Inches(3.0),
'Georgia', 54, hl_color)

# Body
if body:
add_textbox(slide, body, LM, Inches(5.5), TW, Inches(4.0),
'Calibri', 28, body_color)

# Italic close
if closing:
add_textbox(slide, closing, LM, Inches(10.5), TW, Inches(2.0),
'Georgia', 39, close_color, italic=True)

return slide
```

def build_slide_3(prs, copy, italic_accent_color, gradient_stops=None):
“””
S3 Reveal — dark background (FGDeep or DarkBrown gradient). Cream text. Teal/Gold italic.
“””
slide = prs.slides.add_slide(prs.slide_layouts[6])

```
label = s(copy.get('label', ''))
headline = s(copy.get('headline', ''))
hl_it = s(copy.get('headline_italic', ''))
body = s(copy.get('body', ''))
closing = s(copy.get('closing', ''))

stops = gradient_stops or GRADIENTS['fg_deep']
add_gradient_bg(slide, stops)

if label:
add_section_label(slide, label, italic_accent_color)

if headline and hl_it:
add_headline_mixed(slide, headline + ' ', hl_it,
LM, Inches(1.6), TW, Inches(3.5),
'Georgia', 54, C['offwhite'], italic_accent_color)
elif headline:
add_textbox(slide, headline, LM, Inches(1.6), TW, Inches(3.0),
'Georgia', 54, C['offwhite'])

if body:
add_textbox(slide, body, LM, Inches(5.5), TW, Inches(4.0),
'Calibri', 28, C['offwhite'])

if closing:
add_textbox(slide, closing, LM, Inches(10.5), TW, Inches(2.0),
'Georgia', 36, italic_accent_color, italic=True)

return slide
```

def build_slide_4(prs, copy, italic_accent_color, bg_color=None, gradient_stops=None):
“””
S4 Reframe — warm color or green-teal gradient. Cream all text.
“””
slide = prs.slides.add_slide(prs.slide_layouts[6])

```
label = s(copy.get('label', ''))
headline = s(copy.get('headline', ''))
hl_it = s(copy.get('headline_italic', ''))
body = s(copy.get('body', ''))
closing = s(copy.get('closing', ''))

if gradient_stops:
add_gradient_bg(slide, gradient_stops)
else:
add_solid_bg(slide, bg_color or C['terra'])

if label:
add_section_label(slide, label, C['offwhite'])

if headline and hl_it:
add_headline_mixed(slide, headline + ' ', hl_it,
LM, Inches(1.6), TW, Inches(4.0),
'Georgia', 69, C['offwhite'], C['offwhite'])
elif headline:
add_textbox(slide, headline, LM, Inches(1.6), TW, Inches(3.5),
'Georgia', 69, C['offwhite'])

if body:
add_textbox(slide, body, LM, Inches(6.0), TW, Inches(3.5),
'Calibri', 28, C['offwhite'])

if closing:
add_textbox(slide, closing, LM, Inches(10.5), TW, Inches(2.0),
'Georgia', 36, C['offwhite'], italic=True)

return slide
```

def build_slide_5(prs, copy, italic_accent_color):
“””
S5 Invite — always Cream.
Setup 48pt. Emphasis 72pt italic accent. Body 26pt. CTA italic. Quiz italic small. Accent mark.
“””
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, C[‘cream’])

```
setup = s(copy.get('setup', ''))
emphasis = s(copy.get('emphasis', ''))
body = s(copy.get('body', ''))
cta = s(copy.get('cta', ''))
quiz = s(copy.get('quiz', ''))

if setup:
add_textbox(slide, setup,
LM, Inches(1.5), TW, Inches(2.0),
'Georgia', 48, C['dark_brown'])

if emphasis:
add_textbox(slide, emphasis,
LM, Inches(4.0), TW, Inches(3.0),
'Georgia', 72, italic_accent_color, italic=True)

if body:
add_textbox(slide, body,
LM, Inches(7.2), TW, Inches(1.8),
'Calibri', 26, C['dark_brown'])

if cta:
add_textbox(slide, cta,
LM, Inches(9.2), TW, Inches(1.5),
'Georgia', 38, italic_accent_color, italic=True)

if quiz:
add_textbox(slide, quiz,
LM, Inches(11.0), TW, Inches(1.0),
'Calibri', 21, italic_accent_color, italic=True)

# Accent mark — always Terracotta on cream background
add_accent_mark(slide, C['terra'])
return slide
```

# ─────────────────────────────────────────────

# MAIN CAROUSEL BUILDER

# ─────────────────────────────────────────────

def build_carousel(data):
“””
Main entry point. Accepts Claude’s JSON dict.
Returns PPTX bytes.
“””
carousel_type = int(data.get(‘carousel_type’, 3))
photo_assignment = data.get(‘photo_assignment’, ‘PHOTO-A’)
italic_accent_color = data.get(‘italic_accent_color’, C[‘teal’])
gradient_spec = str(data.get(‘gradient_spec’, ‘fg_deep’) or ‘fg_deep’)
layout_notes = str(data.get(‘layout_notes’, ‘’) or ‘’)

```
s1 = data.get('slide_1_copy', {})
s2 = data.get('slide_2_copy', {})
s3 = data.get('slide_3_copy', {})
s4 = data.get('slide_4_copy', {})
s5 = data.get('slide_5_copy', {})

# Determine S4 background from gradient_spec or layout_notes
s4_gradient = None
s4_solid = C['terra']
if 'green_teal' in gradient_spec or 'green' in layout_notes.lower():
s4_gradient = GRADIENTS['green_teal']
s4_solid = None
elif 'terra' in gradient_spec:
s4_solid = C['terra']

# Determine S3 gradient
s3_gradient = GRADIENTS['fg_deep']
if 'sage' in gradient_spec.lower():
s3_gradient = GRADIENTS['fg_deep']

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Build all 5 slides
build_slide_1(prs, s1, carousel_type, photo_assignment,
italic_accent_color, gradient_spec)

# S2: cream for Type 3, gradient for Type 2
s2_gradient = None
if carousel_type == 2:
s2_gradient = GRADIENTS.get('terra_salmon')
build_slide_2(prs, s2, italic_accent_color,
is_dark_bg=False, gradient_stops=s2_gradient)

build_slide_3(prs, s3, italic_accent_color, gradient_stops=s3_gradient)

build_slide_4(prs, s4, italic_accent_color,
bg_color=s4_solid, gradient_stops=s4_gradient)

build_slide_5(prs, s5, italic_accent_color)

# Save to bytes
buf = io.BytesIO()
prs.save(buf)
buf.seek(0)
return buf
```

# ─────────────────────────────────────────────

# QUALITY GATES

# ─────────────────────────────────────────────

def run_quality_gates(data):
“””
Check Claude’s output before building.
Returns list of violations. Empty list = pass.
“””
violations = []
all_text = json.dumps(data)

```
if '\u2014' in all_text:
violations.append('Em dash found in output')
if 'glimpsing' in all_text.lower():
violations.append('"glimpsing" found in output')
if 'and it matters' in all_text.lower():
violations.append('"and it matters" found in output')
if 'there\'s a difference' in all_text.lower():
violations.append('"there\'s a difference" found in output')

hashtags = data.get('hashtags', [])
if len(hashtags) != 4:
violations.append(f'Expected 4 hashtags, got {len(hashtags)}')

caption = str(data.get('caption', '') or '')
if ' she ' in caption.lower() or ' her ' in caption.lower():
violations.append('Caption uses she/her for reader — should use you/your')

s1 = data.get('slide_1_copy', {})
if s1.get('label'):
violations.append('Slide 1 has a section label — must be removed')

return violations
```

# ─────────────────────────────────────────────

# FLASK ROUTES

# ─────────────────────────────────────────────

@app.route(’/build-carousel’, methods=[‘POST’])
def build_carousel_endpoint():
“””
POST /build-carousel
Body: Claude’s JSON output
Returns: PPTX file binary

```
Make should:
1. POST the Claude JSON response to this endpoint
2. Store the returned binary as a .pptx file
3. Attach it to the Gmail approval email
"""
# Accept either:
# 1. A proper JSON object (Content-Type: application/json)
# 2. A raw text string containing JSON (what Make sends from Claude's text output)
data = None
raw = request.get_data(as_text=True).strip()

# Try parsing as JSON directly first
try:
data = json.loads(raw)
except Exception:
pass

# If that failed, try to extract JSON from within the text
# Claude sometimes wraps JSON in extra text or backticks
if data is None:
import re
match = re.search(r'\{.*\}', raw, re.DOTALL)
if match:
try:
data = json.loads(match.group())
except Exception:
pass

if data is None:
return jsonify({'error': f'Could not parse JSON from request body. Raw content starts with: {raw[:200]}'}), 400

# Run quality gates
violations = run_quality_gates(data)
if violations:
return jsonify({
'error': 'Quality gate failures — do not deliver',
'violations': violations
}), 422

try:
pptx_buf = build_carousel(data)
topic = data.get('topic', 'carousel').replace(' ', '_')[:40]
filename = f'RELEASED_{topic}.pptx'
return send_file(
pptx_buf,
mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
as_attachment=True,
download_name=filename
)
except Exception as e:
return jsonify({'error': f'Build failed: {str(e)}'}), 500
```

@app.route(’/health’, methods=[‘GET’])
def health():
return jsonify({‘status’: ‘ok’, ‘service’: ‘RELEASED Carousel Builder’})

@app.route(’/test’, methods=[‘GET’])
def test_build():
“””
GET /test — builds a sample carousel using dummy copy.
Use to verify the service is working after deployment.
“””
sample = {
“topic”: “Breaking the Fear-Symptom Loop”,
“carousel_type”: 3,
“photo_assignment”: “PHOTO-A”,
“italic_accent_color”: “#6fafb8”,
“gradient_spec”: “fg_deep”,
“layout_notes”: “green-teal gradient on S4”,
“slide_1_copy”: {
“headline”: “The more you focus on a symptom,”,
“headline_italic”: “the worse it gets.”,
“sub_statement”: “You’ve probably noticed that. Here’s why it’s happening.”
},
“slide_2_copy”: {
“label”: “THE LOOP IN ACTION”,
“headline”: “Your nervous system learned to”,
“headline_italic”: “scan.”,
“body”: “When something felt threatening, your body started watching for it. That made sense then. The problem is it never got the signal that the threat had passed.”,
“closing”: “So it’s still watching. Every day.”
},
“slide_3_copy”: {
“label”: “WHAT’S ACTUALLY HAPPENING”,
“headline”: “This is the fear-symptom”,
“headline_italic”: “loop.”,
“body”: “Fear activates the nervous system. The nervous system produces symptoms. The symptoms create more fear. Nothing in that cycle tells your body it’s safe.”,
“closing”: “The loop doesn’t break through willpower.”
},
“slide_4_copy”: {
“label”: “THE WAY OUT”,
“headline”: “You don’t think your way out of”,
“headline_italic”: “a nervous system pattern.”,
“body”: “The body has to receive the signal directly. That’s what this healing is for.”,
“closing”: “A gap can always be filled.”
},
“slide_5_copy”: {
“setup”: “If you’ve been trying to manage this on your own…”,
“emphasis”: “That’s not a character failure.”,
“body”: “You’ve been without support that goes as deep as the pattern does.”,
“cta”: “The stress bucket quiz in my bio is a good place to start.”,
“quiz”: “Link’s there.”
},
“section_labels”: [
“THE LOOP IN ACTION”,
“WHAT’S ACTUALLY HAPPENING”,
“THE WAY OUT”
],
“caption”: “Your nervous system isn’t broken. It learned something, and it held on. That’s what bodies do. The question isn’t what’s wrong with you. It’s what happened, and whether anyone has helped your body understand it’s over.\n\nIf that’s you, the quiz in my bio is a starting point.”,
“hashtags”: [”#nervoussystemhealing”, “#somatichealing”, “#traumarecovery”, “#highlysensitiveperson”],
“competitor_post”: {
“account”: “@example_account”,
“why_it_fits”: “High save rate on a somatic awareness post targeting similar audience.”,
“preview”: “Your body isn’t the enemy. It’s just been trying to keep you safe…”,
“source_url”: “https://instagram.com/p/example”
}
}

```
violations = run_quality_gates(sample)
if violations:
return jsonify({'violations': violations}), 422

pptx_buf = build_carousel(sample)
return send_file(
pptx_buf,
mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
as_attachment=True,
download_name='RELEASED_test_carousel.pptx'
)
```

if **name** == ‘**main**’:
port = int(os.environ.get(‘PORT’, 5000))
app.run(host=‘0.0.0.0’, port=port, debug=False)
